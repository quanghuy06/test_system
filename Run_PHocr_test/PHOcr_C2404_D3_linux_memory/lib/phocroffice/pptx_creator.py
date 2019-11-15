# -*- encoding: utf-8 -*-
# Copyright (C) 2017 TSDV TTEC.  All rights reserved.

from __future__ import division

import warnings
warnings.filterwarnings('ignore', category=DeprecationWarning)

import sys
import os

from pptx import Presentation
from pptx.table_processing import TableProcessing
from pptx.enum.table_style import TABLE_STYLE
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.oxml.simpletypes import ST_TextDirectionType
from phocr_shared.shared import Pt, Pixel, Inches

from phocr_elements.office_enum import PageDirection, LineDirection
from phocr_elements.office_object_tag import OfficeObjectTag
from phocr_shared.phocr_config import PHOcrConstant
from phocr_shared.phocr_common import Alignment, \
    PHOcrCommon, TextDirection, DocumentType
from phocr_shared.bullets_and_numberings import update_bullets_and_numbering, \
    BulletNumberingContainer
from phocr_shared.alignment import string_to_alignment_pptx, analysis_indent, \
    calculate_left_indent, indent_one_line_paragraph
from phocr_shared import text_box, ocr_box
from phocr_shared.alignment_enum import AlignmentEnum
from phocr_shared.box_container import BoxContainer
from phocr_shared.text_direction import WritingDirection, TextlineOrder, \
    OrientationTransformer
from phocr_shared.logger import OCRLogLevel
from base_office_creator import BaseOfficeCreator
from phocr_shared.word_format import WordFormat
from phocr_shared.phocr_dpi import PHOcrDPI


class PptxCreator(BaseOfficeCreator):
    """Creator to generate .pptx file"""

    # White space width is equal 1/4 portion of font size
    WHITE_SPACE_RATIO = 0.25

    def __init__(self, filename, language, debug, working_directory):
        super(PptxCreator, self).__init__(language, working_directory, debug)
        self.document_type = DocumentType.PPTX
        self.output_file = filename + '.pptx'
        self.text_direction = TextDirection.LRTB
        self.width = 0
        self.height = 0
        self.x_height = 1
        self.page_deskew_angle = 0
        self.page_orientation = 0
        self.page_writing_direction = 0
        self.page_textline_order = 2
        self.dpi = PHOcrDPI(PHOcrConstant.DEFAULT_DPI, PHOcrConstant.DEFAULT_DPI)
        self.alpha_number_refer_uper = [
            'A', 'B', 'C', 'D',
            'E', 'F', 'G', 'H',
            'I', 'K', 'L'
        ]
        self.alpha_number_refer_lower = [
            'a', 'b', 'c', 'd',
            'e', 'f', 'g', 'h',
            'i', 'k', 'l'
        ]
        self.bullet_numberings = BulletNumberingContainer()
        self.page = {}

        # Direction
        self.transformer = OrientationTransformer()

        self.table_processing = TableProcessing()

    def set_width_and_height_for_column(self, element, table):
        """
        This function will generate table from xml information

        Parameters
        ----------
        element : Object
            XML object of table
        table : Table
            Table object of pptx library

        Returns
        -------

        """
        # Get width of columns, height of rows in table
        col_width_arr, col_height_arr = PHOcrCommon.get_columns_size(
            element
        )
        # Set width of columns
        ci = 0
        for ecw in col_width_arr:
            table.columns[ci].width = Pixel(ecw, self.dpi.horizontal_resolution)
            ci += 1
        # Set height of rows
        ri = 0
        for erh in col_height_arr:
            table.rows[ri].height = Pixel(erh, self.dpi.vertical_resolution)
            ri += 1

    def generate_data_in_cell(
            self, slide, tf, cell, column, page_box, para_index):
        for sub_photo in column.photos:
            self.generate_photo(slide, sub_photo)
        for sub_carea in column.careas:
            cell_box = ocr_box.OCRBox.from_xml_element(
                column,
                self.dpi,
                Pixel
            )
            box_container = BoxContainer(
                box=cell_box,
                parent_box=page_box
            )

            for paragraph in sub_carea.paragraphs:
                is_first_para = False
                if para_index == 0:
                    is_first_para = True
                p = PptxCreator.get_next_text_frame_paragraph(
                    first_paragraph_in_text_frame=is_first_para,
                    text_frame=tf
                )

                self.add_paragraph_to_carea(p,
                                            sub_carea,
                                            paragraph,
                                            box_container)
                para_index += 1
        self.table_processing.set_fill_for_cell(
            cell=cell,
            missing_edge_xml=column
        )

    def merge_row_and_cell(self, table, column, row_id, column_id):
        col_span = column.col_span
        row_span = column.row_span
        # If both col_span and row_span > 1
        if col_span > 1 and row_span > 1:
            # Merge cells vertically and horizontally
            self.table_processing.merge_cells(
                table=table,
                start_row_idx=row_id - 1,
                end_row_idx=row_id - 1 + row_span - 1,
                start_col_idx=column_id - 1,
                end_col_idx=column_id - 1 + col_span - 1)
        elif col_span > 1:  # If horizontal merge
            # Merge cells horizontally
            self.table_processing.merge_horizontally_cells(
                table=table,
                row_idx=row_id - 1,
                start_col_idx=column_id - 1,
                end_col_idx=column_id - 1 + col_span - 1)
        elif row_span > 1:  # If vertical merge
            # Merge cells vertically
            self.table_processing.merge_vertically_cells(
                table=table,
                start_row_idx=row_id - 1,
                end_row_idx=row_id - 1 + row_span - 1,
                col_idx=column_id - 1)

    def set_alignment_in_cell(self, cell, column):
        num_col_elements = column.length()
        if num_col_elements > 0 and \
                not (num_col_elements == 1 and len(column.photos) == 1):
            top_distance, \
            _left_distance, \
            bottom_distance, \
            _right_distance, \
            _char_width = self.calculate_alignment_properties(column)

            # Get left-right-top-bottom distance
            top_val = abs(top_distance - column.y)
            bottom_val = abs(column.y + column.h - bottom_distance)

            # Set vertical alignment of cell
            vert_alignment = Alignment.TOP
            if abs(top_val - bottom_val) < 60:
                vert_alignment = Alignment.MIDDLE
            elif bottom_val != 0:
                vert_ratio = round(float(top_val) / bottom_val, 3)
                if abs(top_val - bottom_val) > 0.8:
                    vert_alignment = Alignment.TOP
                elif vert_ratio > 1.2:
                    vert_alignment = Alignment.BOTTOM
                else:
                    vert_alignment = Alignment.MIDDLE
            if vert_alignment == Alignment.MIDDLE:  # If middle
                cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            elif vert_alignment == Alignment.BOTTOM:  # If bottom
                cell.vertical_anchor = MSO_VERTICAL_ANCHOR.BOTTOM

    @classmethod
    def set_margin_in_cell(cls, cell):
        emu_lr_margin = 0  # 0 instead of 91440 (0.10 inches)
        emu_tb_margin = 0  # 0 instead of 45720 (0.05 inches)
        cell.margin_left = emu_lr_margin
        cell.margin_right = emu_lr_margin
        cell.margin_top = emu_tb_margin
        cell.margin_bottom = emu_tb_margin

    def generate_data_in_table(self, slide, table, element):
        page_box = ocr_box.OCRBox.from_xml_element(self.page, self.dpi, Pixel)
        for row in element.rows:
            for column in row.cells:
                # Write body cells. Cell is a table cell
                cell = table.cell(row.id - 1, column.id - 1)

                self.merge_row_and_cell(table, column, row.id, column.id)
                tf = cell.text_frame
                # Check style of cell
                self.set_alignment_in_cell(cell, column)

                # We'll reset margin of cell like Word
                PptxCreator.set_margin_in_cell(cell)

                # Print sub element
                para_index = 0
                self.generate_data_in_cell(
                    slide, tf, cell, column, page_box, para_index)

                PptxCreator.fill_cell_background_color(cell, column)
                # In my knowledge, we can't insert a table inside cell
                # So in case structure of cell has table,
                # for now, we'll ingore it
                # In the future, if we have other way to generate table
                # in cell, we can add more code here
                # if sub_element.tag == 'ocr_table_mask':
                # self.generate_table(slide, sub_element)

    def generate_table(self, slide, element):
        """
        This function will generate table from xml information

        Parameters
        ----------
        slide : Slide
            Instance of Slides object
        element : Object
            XML object of table

        Returns
        -------

        """

        table_pos = PHOcrCommon.get_object_position(element)
        cx = Pixel(table_pos.x, self.dpi.horizontal_resolution)
        cy = Pixel(table_pos.y, self.dpi.vertical_resolution)
        cw = Pixel(table_pos.w, self.dpi.horizontal_resolution)
        ch = Pixel(table_pos.h, self.dpi.vertical_resolution)
        tableStyleId = TABLE_STYLE.NoStyle_TableGrid
        shapes = slide.shapes
        table = shapes.add_table(
            element.num_rows, element.num_cells,
            cx, cy, cw, ch, tableStyleId
        ).table
        self.table_processing.clear_table_default(table,
                                                  element.num_rows,
                                                  element.num_cells)
        self.generate_data_in_table(slide, table, element)
        self.set_width_and_height_for_column(element, table)

    @classmethod
    def get_next_text_frame_paragraph(
            cls,
            first_paragraph_in_text_frame,
            text_frame):
        if first_paragraph_in_text_frame:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.clear()
        return p

    @classmethod
    def fill_cell_background_color(cls, cell, cell_xml):
        if cell_xml.highlight_color is not None:
            cell.fill.solid()
            cell.fill.fore_color.rgb = PHOcrCommon.get_color(
                cell_xml.highlight_color
            )
            return

        if len(cell_xml.careas) != 1:
            return

        background_color = cell_xml.careas[0].background_color
        if background_color is not None:
            # cell is a table cell
            # set fill type to solid color first
            cell.fill.solid()
            # set foreground (fill) color to a specific RGB color
            cell.fill.fore_color.rgb = PHOcrCommon.get_color(background_color)

    @staticmethod
    def calculate_average_char_width_in_line(line_xml):
        word_w = 0
        num_chars = 0
        for word_xml in line_xml.words:
            word_w += word_xml.w
            num_chars += len(word_xml.characters)
        return word_w, num_chars

    def calculate_average_char_width_in_carea(self, carea_xml):
        word_w = 0
        num_chars = 0
        for par_xml in carea_xml.paragraphs:
            for line_xml in par_xml.lines:
                word_w_tmp, num_chars_tmp = \
                    self.calculate_average_char_width_in_line(line_xml)
                word_w += word_w_tmp
                num_chars += num_chars_tmp

        if num_chars > 0:
            char_w = round(float(word_w) / num_chars, 3)
            return char_w
        else:
            return -1

    def calculate_alignment_properties(self, column):
        top_distance = sys.maxsize
        left_distance = sys.maxsize
        bottom_distance = 0
        right_distance = 0
        char_width = 0
        num_careas = 0
        for sub_carea in column.careas:
            num_careas += 1
            sub_element_pos = PHOcrCommon.get_object_position(sub_carea)
            if sub_element_pos.x < left_distance:
                left_distance = sub_element_pos.x
            if sub_element_pos.y < top_distance:
                top_distance = sub_element_pos.y
            if sub_element_pos.x + sub_element_pos.w > right_distance:
                right_distance = sub_element_pos.x + sub_element_pos.w
            if sub_element_pos.y + sub_element_pos.h > bottom_distance:
                bottom_distance = sub_element_pos.y + sub_element_pos.h
            char_width += self.calculate_average_char_width_in_carea(
                sub_carea
            )

        if num_careas > 0:
            char_width = char_width / num_careas

        return (top_distance,
                left_distance,
                bottom_distance,
                right_distance,
                char_width)

    def generate_photo(self, slide, photo):
        """
        Generate a photo into MS Power Point Document representation

        Parameters
        ----------
        slide
        photo

        Returns
        -------

        """
        photo_pos = PHOcrCommon.get_object_position(photo)
        # Only add photo if it is bigger than threshold
        if PHOcrCommon.is_photo_ok(photo_pos.h, photo_pos.w, self.x_height):
            cx = Pixel(photo_pos.x, self.dpi.horizontal_resolution)
            cy = Pixel(photo_pos.y, self.dpi.vertical_resolution)
            cw = Pixel(photo_pos.w, self.dpi.horizontal_resolution)
            ch = Pixel(photo_pos.h, self.dpi.vertical_resolution)
            absolute_photo_path = os.path.join(self.working_directory, photo.path)
            _ = slide.shapes.add_picture(absolute_photo_path, cx, cy, cw, ch)

    def get_space_text(self, p, para, box_container, paragraph_information,
                       word, first_word):
        spaces_before = word.spaces_before
        if spaces_before is None:
            spaces_before = 0
        # Assign value for run_space
        if spaces_before == 0 and not first_word:
            # if spaces_before not defined, it should be 1 space
            # however, when it is first word, don't need space anymore
            spaces_before = 1

        # Add space
        if spaces_before > 3 and AlignmentEnum.is_left(para.alignment):
            # Only when the number of space large enough
            tab_stops = p.tab_stops
            tab_position = Pixel(word.x - box_container.left_margin(), self.dpi.horizontal_resolution)

            # If the margin of page left are OK, we will add the tab stop
            tab = tab_stops.add_tab(tab_position)
            space_text = '\t'

            paragraph_information['last_tab'] = tab
            paragraph_information['last_tab_word'] = word
        else:
            space_text = ' ' * spaces_before
        return space_text

    def set_word_style(self, run, word):
        font = run.font
        if word.font and len(word.font) > 0:
            font.name = word.font
        if word.size and len(word.size) > 0 and float(word.size):
            font.size = Pt(float(word.size))
        if word.bold:
            font.bold = True
        if word.italic:
            font.italic = True
        if word.underline:
            font.underline = True
        if word.color is not None:
            run.font.color.rgb = PHOcrCommon.get_color(word.color)

    def generate_word(
            self,
            p,
            para,
            word,
            box_container,
            paragraph_information,
            index=1):
        """
        Generate a word into MS Power Point Document representation

        Parameters
        ----------
        p
        para
        word
        box_container
        paragraph_information
        index

        Returns
        -------

        """
        first_word = index == 1
        # Add a run for space
        run_space = p.add_run() if not first_word else None
        # Add a run for text
        run = p.add_run()
        font = run.font
        self.set_word_style(run, word)

        space_text = self.get_space_text(
            p, para, box_container, paragraph_information,
            word, first_word
        )
        self.setup_space_run(run_space, font, space_text)

        run.text = word.value
        return run

    def generate_line(
            self,
            p,
            para,
            line,
            line_position,
            box_container,
            paragraph_information,
            is_first_line=False,
            is_list_type=False):
        """
        Generate a line into MS Power Point Document representation

        Parameters
        ----------
        p
        para
        line
        box_container
        paragraph_information
        is_first_line
        is_list_type

        Returns
        -------

        """
        index = 1
        is_last_line_of_paragraph = line == para.lines[-1]
        for word_position, word in enumerate(line.words):
            is_last_word_of_line = word == line.words[-1]
            ignore = False
            # We won't generate first word of paragraph classified as
            # bullets and numbering.
            # that word is symbol of bullets or numbering.
            # We generate it in another place
            if is_first_line and is_list_type and index == 1:
                ignore = True
            if not ignore:
                # Finding previous word
                previous_word = None
                if word_position - 1 >= 0:
                    previous_word = line.words[word_position - 1]
                if word_position == 0 and line_position - 1 >= 0:
                    previous_word = para.lines[line_position - 1].words[-1]

                # Checking word and previous word are same format
                word_format = WordFormat(word)
                previous_word_format = WordFormat(previous_word)
                is_samle_format = word_format.is_same_format(
                    previous_word_format
                )
                text_run = None
                if is_samle_format and len(p.runs) > 0:
                    text_run = p.runs[-1]
                    first_word = index == 1
                    space_text = self.get_space_text(
                        p, para, box_container, paragraph_information,
                        word, first_word
                    )
                    text_run.text = text_run.text + space_text + word.value
                else:
                    text_run = self.generate_word(
                        p,
                        para,
                        word,
                        box_container,
                        paragraph_information,
                        index
                    )
                if is_last_word_of_line and not is_last_line_of_paragraph:
                    text_run.text += ' '
            index += 1

    def generate_paragraph(
            self,
            p,
            paragraph,
            box_container,
            paragraph_information,
            is_list_type=False):
        """
        Generate a paragraph into MS Power Point Document representation

        Parameters
        ----------
        p
        paragraph
        box_container
        paragraph_information
        is_list_type

        Returns
        -------

        """
        # Set style for paragraph
        alignment = PHOcrCommon.get_paragraph_alignment(
            paragraph,
            self.text_direction,
            self.x_height
        )
        p.alignment = string_to_alignment_pptx(alignment)

        for index, line in enumerate(paragraph.lines):
            is_first_line = index == 0
            self.generate_line(p, paragraph, line,
                               index,
                               box_container,
                               paragraph_information,
                               is_first_line, is_list_type)

        if len(paragraph.lines) > 1:
            line_spacing_pix = PHOcrCommon.calculate_line_spacing(
                paragraph,
                self.text_direction,
                self.document_type
            )

            # Initialize default value of the right resolution
            resolution = PHOcrConstant.DEFAULT_DPI
            if self.text_direction == TextDirection.LRTB:
                resolution = self.dpi.horizontal_resolution
            elif self.text_direction == TextDirection.TBRL:
                resolution = self.dpi.vertical_resolution
            line_spacing_inches = line_spacing_pix / resolution
            if line_spacing_inches < 0.005:
                line_spacing_inches = 0.005
            if line_spacing_inches > 0:
                p.line_spacing = Inches(line_spacing_inches)

    def get_text_frame(self, carea, textbox):
        writing_direction = carea.writing_direction
        textline_order = carea.textline_order
        tf = textbox.text_frame
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        text_direction = ST_TextDirectionType.HORZ
        # If direction of layout is Japanese or Chinese
        if writing_direction == PageDirection.VERTICAL and \
                textline_order == LineDirection.RIGHT_TO_LEFT:
            text_direction = ST_TextDirectionType.EAVERT
        else:
            check_direction = False
            for paragraph in carea.paragraphs:
                for line in paragraph.lines:
                    text_angle = int(line.text_angle)
                    if text_angle == 90:
                        text_direction = ST_TextDirectionType.VERT
                    elif text_angle == 270:
                        text_direction = ST_TextDirectionType.VERT270
                    if text_angle != -1:
                        check_direction = True
                        break
                if check_direction:
                    break
        tf.text_direction = text_direction
        return tf

    def generate_carea(self, slide, carea):
        """
        Generate a content area into MS Power Point Document representation

        Parameters
        ----------
        slide
        carea

        Returns
        -------

        """
        if len(carea.paragraphs) == 0:
            return

        # NamLD: because carea with one line usually meet error
        # from font size/name
        # Therefore which carea have one line will be break to two line
        # if the size of box are exactly same with image.
        # That carea will overlap other
        # This code segment will increasing the width of carea
        # if it contain only one line
        # Carea have 1 paragraph, that paragraph have 1 line

        carea_x, carea_y, carea_w, carea_h = \
            text_box.calculate_text_box_size(carea)
        carea_container_box = ocr_box.OCRBox(self.dpi,
                                             carea_x, carea_y,
                                             carea_w, carea_h,
                                             Pixel)

        page_box = ocr_box.OCRBox.from_xml_element(self.page, self.dpi, Pixel)
        box_container = \
            BoxContainer(box=carea_container_box, parent_box=page_box)

        # In PPTX, background color of carea will be background color of whole
        # text box
        bg_color = None
        if carea.background_color is not None:
            bg_color = PHOcrCommon.get_color(carea.background_color)
        cx = Pixel(carea_x, self.dpi.horizontal_resolution)
        cy = Pixel(carea_y, self.dpi.vertical_resolution)
        cw = Pixel(carea_w, self.dpi.horizontal_resolution)
        ch = Pixel(carea_h, self.dpi.vertical_resolution)
        textbox = slide.shapes.add_textbox(cx, cy, cw, ch, bg_color)
        tf = self.get_text_frame(carea, textbox)

        for index, paragraph in enumerate(carea.paragraphs):
            p = PptxCreator.get_next_text_frame_paragraph(
                first_paragraph_in_text_frame=(index == 0),
                text_frame=tf
            )

            self.add_paragraph_to_carea(p, carea, paragraph, box_container)
        tf.word_wrap = True

    def set_style_for_paragraph(self, p, carea, paragraph, box_container):
        """
        Set style information of paragraph as first_line_indent, alignment, ...

        Parameters
        ----------
        p : _Paragraph
            Instance of Paragraph in pptx library XML object of page
        carea : Object
            XML object of carea which contains paragraphs
        paragraph : Object
            XML object of paragraph inside carea
        box_container : BoxContainer
            Contains relative information of object

        Returns
        -------

        """

        # Calculate alignment
        p.alignment = string_to_alignment_pptx(paragraph.alignment)
        (left_indent, right_indent, first_line_indent) = \
            analysis_indent(
                p, paragraph, self.dpi, Pixel, block=carea,
                page={
                    'margin_left':
                        Pixel(box_container.left_margin(),
                              self.dpi.horizontal_resolution).emu,
                    'margin_right':
                        Pixel(box_container.right_margin(),
                              self.dpi.horizontal_resolution).emu,
                    'margin_top':
                        Pixel(box_container.top_margin(),
                              self.dpi.vertical_resolution).emu,
                    'margin_bottom':
                        Pixel(box_container.bottom_margin(),
                              self.dpi.vertical_resolution).emu,
                    'width':
                        Pixel(box_container.parent_w_px,
                              self.dpi.horizontal_resolution).emu,
                    'height':
                        Pixel(box_container.parent_h_px,
                              self.dpi.vertical_resolution).emu,
                    'text_line_order':
                        self.transformer.text_line_order,
                    'writing_direction':
                        self.transformer.writing_direction,
                    'orientation':
                        self.transformer.orientation
                }
            )

        # Correct left, right indent with one line paragraph
        left_indent, right_indent = indent_one_line_paragraph(
            paragraph,
            left_indent,
            right_indent,
            self.transformer,
            Pixel
        )

        # first_line_indent is comparing with paragraph alignment,
        # however returned first_line_indent is
        p.first_line_indent = first_line_indent
        p.left_margin = max(calculate_left_indent(p, left_indent), 0)
        p.word_wrap = True

        if AlignmentEnum.is_right(paragraph.alignment) and \
                len(paragraph.lines) == 1:
            # If right alignment and and have one line, adding space to after,
            # because pptx don't have right indent
            one_space_size = Pixel(13, 300).emu  # By measurement
            num_spaces = int(right_indent / one_space_size)

            if num_spaces > 3:
                space_after = p.add_run()
                space_after.text = ' ' * num_spaces
                space_after.font.name = 'Arial'
                # by measure one_space size above
                space_after.font.size = Pt(11.0)

                # Because pptx only care to printed character. Therefore, we've
                # the type and trick here. Adding . with font size very small
                very_small = p.add_run()
                very_small.text = '.'
                very_small.font.name = 'Arial'
                very_small.font.size = Pt(1.0)

    def add_paragraph_to_carea(self, p, carea, paragraph, box_container):
        """
        Add paragraph to carea
        - Add list for paragraph
        - Generate text for paragraph
        - Calculate alignment for paragraph

        Parameters
        ----------
        p
        carea
        paragraph
        box_container

        Returns
        -------

        """
        # List type
        is_list_type = False
        if paragraph.is_list and \
                len(paragraph.lines) > 0 and len(paragraph.lines[0].words) > 1:
            is_list_type = update_bullets_and_numbering(self,
                                                        None,
                                                        p,
                                                        paragraph,
                                                        carea.x,
                                                        Pixel,
                                                        box_container)
        paragraph_information = {
            'last_tab': None,
            'last_tab_word': None
        }
        # Generate paragraph
        self.generate_paragraph(p, paragraph,
                                box_container,
                                paragraph_information,
                                is_list_type)
        self.set_style_for_paragraph(p, carea, paragraph, box_container)
        self.correct_tab_stop_when_too_tight(
            paragraph,
            box_container,
            paragraph_information)

    def get_page_information(self, page):
        """
        Get information of page and store

        Parameters
        ----------
        page : Object
            XML object of page

        Returns
        -------

        """

        # Get dpi of image
        self.dpi = PHOcrCommon.get_dpi(page.dpi)
        if page.x_height > 1:
            self.x_height = page.x_height

        self.page_deskew_angle = float(page.deskew_angle)
        self.page_orientation = page.orientation
        self.page_writing_direction = page.writing_direction
        self.page_textline_order = page.textline_order
        self.width = page.w
        self.height = page.h

        self.transformer.orientation = page.orientation
        self.transformer.writing_direction = page.writing_direction
        self.transformer.text_line_order = page.textline_order

        self.text_direction = TextDirection.LRTB
        # If we found layout is Chinese or Japanese
        if self.page_writing_direction == \
                WritingDirection.WRITING_DIRECTION_TOP_TO_BOTTOM and \
                self.page_textline_order == \
                TextlineOrder.TEXTLINE_ORDER_RIGHT_TO_LEFT:
            self.text_direction = TextDirection.TBRL

    def create_new_slide_for_page(self, pptx, page):
        """
        Create new slide based on xml information of page

        Parameters
        ----------
        pptx : Presentation
            Instance of pptx library for call APIs
        page : Object
            XML object of page

        Returns
        -------

        """

        # Create blank slide
        blank_slide_layout = pptx.slide_layouts[6]
        slide = pptx.slides.add_slide(blank_slide_layout)

        # Photos are generated firstly to make sure that photo behind text
        for photo in page.photos:
            self.generate_photo(slide, photo)
        for element in page.elements:
            if element.tag == OfficeObjectTag.CAREA:
                self.generate_carea(slide, element)
            if element.tag == OfficeObjectTag.TABLE:
                self.generate_table(slide, element)

    def generate_document(self, office_obj):
        """
        Generate a power point document representation

        Parameters
        ----------
        office_obj : Object
            XML object which is parsed from .xml file

        Returns
        -------

        """
        # Create a MS Power Point document
        pptx = Presentation()
        # Information in "properties"
        """
        Expected result for all formats:
        - last modified: (blank)
        - created: (current date with format yyyy-mm-dd)
        - author: "None"
        - last modified by: "Not saved yet"
        """
        core_properties = pptx.core_properties
        core_properties.remove_modified_time()
        from datetime import datetime
        core_properties.created = datetime.utcnow()
        core_properties.author = ''
        core_properties.comments = ''
        core_properties.last_modified_by = ''
        try:
            # Set page layout. Need to covert to English Metric Units (EMU).
            max_width, max_height = office_obj.get_max_size_of_pages()
            pptx.slide_width = Inches(max_width)
            pptx.slide_height = Inches(max_height)

            # Generate slides
            for page in office_obj.pages:
                self.page = page

                # Get information of page
                self.get_page_information(page)

                # Create new slide for page
                self.create_new_slide_for_page(pptx, page)

            # Save the file
            pptx.save(self.output_file)
        except Exception as ex:
            # Logging error
            self.error_logger.log(
                OCRLogLevel.LOG_LEVEL_ERROR,
                'Error occurred. Arguments {0}.'.format(ex.args)
            )
            raise

    def setup_space_run(self, run_space, font, space_text):
        """
        Setup the style of space run

        Parameters
        ----------
        run_space
        font
        space_text

        Returns
        -------

        """
        if run_space is not None:
            space_font = run_space.font
            space_font.name = font.name
            space_font.size = font.size
            run_space.text = space_text

    def correct_tab_stop_when_too_tight(self, paragraph,
                                        box_container,
                                        paragraph_information):
        """
        When box container too much tight, we need reduce the tab stop
        That occurs only with one line paragraph

        Parameters
        ----------
        paragraph
        box_container
        paragraph_information

        Returns
        -------

        """
        if len(paragraph.lines) > 1:
            return

        last_tab_position = paragraph_information['last_tab']
        last_tab_word = paragraph_information['last_tab_word']
        if last_tab_position is None or last_tab_word is None:
            return

        word_box = ocr_box.OCRBox.from_xml_element(
            last_tab_word,
            self.dpi,
            Pixel
        )
        paragraph_box = ocr_box.OCRBox.from_xml_element(
            paragraph,
            self.dpi,
            Pixel
        )
        if paragraph_box.w / box_container.w > 0.95:
            reduce_tab = int((paragraph_box.right - word_box.x) * 0.15)
            last_tab_position.pos = last_tab_position.pos - reduce_tab
